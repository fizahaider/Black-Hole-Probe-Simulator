import io
import re
from typing import List, Dict, Any, Generator
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from dataclasses import dataclass

@dataclass
class Chunk:
    text: str
    metadata: Dict[str, Any]

class SemanticChunker:
    """
    Chunks text by respecting structural boundaries (paragraphs, headers).
    Merges small blocks until target size is reached.
    """
    def __init__(self, target_tokens=400, overlap_tokens=50, max_tokens=1000):
                                  
        self.target_chars = target_tokens * 4
        self.overlap_chars = overlap_tokens * 4
        self.max_chars = max_tokens * 4

    def chunk_text(self, text: str, base_metadata: Dict[str, Any]) -> List[Chunk]:
                                                                                          
        blocks = re.split(r'\n\s*\n', text)
        
        chunks = []
        current_block = ""
        current_meta = base_metadata.copy()
        
        for block in blocks:
            block = block.strip()
            if not block:
                continue
                
                                                                     
            if len(current_block) + len(block) > self.target_chars:
                if current_block:
                    chunks.append(Chunk(text=current_block, metadata=current_meta.copy()))
                                                       
                                                                 
                                                                         
                    current_block = ""
                
                                                                                         
                if len(block) > self.max_chars:
                    sub_chunks = self._recursive_split(block, current_meta)
                    chunks.extend(sub_chunks)
                else:
                    current_block = block
            else:
                if current_block:
                    current_block += "\n\n" + block
                else:
                    current_block = block
        
        if current_block:
             chunks.append(Chunk(text=current_block, metadata=current_meta.copy()))
             
        return chunks

    def _recursive_split(self, text: str, metadata: Dict[str, Any]) -> List[Chunk]:
        """Fallback for massive blocks (like strict legal text without breaks)"""
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.target_chars,
            chunk_overlap=self.overlap_chars,
            length_function=len
        )
        docs = splitter.create_documents([text], metadatas=[metadata])
        return [Chunk(text=d.page_content, metadata=d.metadata) for d in docs]


class StreamingPDFLoader:
    """
    Loads PDF page-by-page to support large files without OOM.
    """
    def __init__(self, file_obj):
        self.file = file_obj

    def load_and_chunk(self, chunker: SemanticChunker, base_metadata: Dict[str, Any]) -> Generator[Chunk, None, None]:
        """
        Yields chunks as we process the PDF.
        """
        if hasattr(self.file, 'seek'):
            self.file.seek(0)
            
        pdf_reader = PyPDF2.PdfReader(self.file)
        
                                                                                               
                                                              
        
        current_section = "Intro"
        
        for i, page in enumerate(pdf_reader.pages):
            text = page.extract_text()
            if not text:
                continue
                
                                              
                                                           
            lines = text.split('\n')
            if lines and len(lines[0]) < 50 and len(lines) > 1:
                potential_header = lines[0].strip()
                                                  
                if not re.search(r'[.,;]$', potential_header):
                    current_section = potential_header
            
            page_meta = base_metadata.copy()
            page_meta.update({
                "page": i + 1,
                "section": current_section
            })
            
                             
            page_chunks = chunker.chunk_text(text, page_meta)
            for chunk in page_chunks:
                yield chunk


class TextLoader:
    """Loads plain text-like files and returns normalized text."""

    def __init__(self, file_obj):
        self.file = file_obj

    def load_text(self) -> str:
        if hasattr(self.file, 'seek'):
            self.file.seek(0)

                                                                         
        raw = self.file.read()
        if isinstance(raw, str):
            return raw
        try:
            return raw.decode('utf-8')
        except UnicodeDecodeError:
            try:
                return raw.decode('utf-16')
            except UnicodeDecodeError:
                return raw.decode('latin-1', errors='replace')


class DOCXLoader:
    """Loads DOCX content and flattens paragraphs to plain text."""

    def __init__(self, file_obj):
        self.file = file_obj

    def load_text(self) -> str:
        if hasattr(self.file, 'seek'):
            self.file.seek(0)
        doc = DocxDocument(self.file)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        return "\n\n".join(paragraphs)


class PPTXLoader:
    """Loads PPTX slides and yields chunked slide text."""

    def __init__(self, file_obj):
        self.file = file_obj

    def load_and_chunk(self, chunker: SemanticChunker, base_metadata: Dict[str, Any]) -> Generator[Chunk, None, None]:
        if hasattr(self.file, 'seek'):
            self.file.seek(0)

        presentation = Presentation(self.file)
        for idx, slide in enumerate(presentation.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    extracted = shape.text.strip()
                    if extracted:
                        text_runs.append(extracted)

            slide_text = "\n\n".join(text_runs).strip()
            if not slide_text:
                continue

            first_line = slide_text.split('\n', 1)[0].strip() or f"Slide {idx + 1}"
            slide_meta = base_metadata.copy()
            slide_meta.update({
                "page": idx + 1,
                "section": first_line[:80],
            })

            for chunk in chunker.chunk_text(slide_text, slide_meta):
                yield chunk
